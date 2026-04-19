from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


class Theme:
    """Centralized visual system for a modern academic presentation."""

    # 16:9 widescreen slide dimensions
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # Color palette
    BG = RGBColor(250, 251, 253)            # soft cool white
    TEXT_PRIMARY = RGBColor(30, 36, 52)     # deep slate
    TEXT_SECONDARY = RGBColor(92, 103, 125) # muted gray-blue
    ACCENT = RGBColor(20, 92, 182)          # restrained academic blue
    ACCENT_SOFT = RGBColor(228, 238, 251)   # subtle tint for cards
    DIVIDER = RGBColor(214, 221, 232)       # neutral divider

    # Typography
    FONT_TITLE = "Calibri"
    FONT_BODY = "Calibri"

    # Spacing system (in inches)
    MARGIN_LEFT = 1.0
    MARGIN_RIGHT = 1.0
    MARGIN_TOP = 0.65
    MARGIN_BOTTOM = 0.45

    GRID_GAP = 0.28


class AcademicPPTTemplate:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Theme.SLIDE_WIDTH
        self.prs.slide_height = Theme.SLIDE_HEIGHT

    # ---------- Utilities ----------

    def _add_background(self, slide):
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = Theme.BG

    def _add_accent_line(self, slide, x, y, w, h=0.05, color=Theme.ACCENT):
        line = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()
        return line

    def _add_textbox(
        self,
        slide,
        text,
        x,
        y,
        w,
        h,
        font_name,
        font_size,
        color,
        bold=False,
        align=PP_ALIGN.LEFT,
        line_spacing=1.2,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.runs[0]
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        return box

    def _add_footer(self, slide, left_text, right_text):
        y = 7.02
        self._add_textbox(
            slide, left_text,
            Theme.MARGIN_LEFT, y, 5.6, 0.26,
            Theme.FONT_BODY, 11, Theme.TEXT_SECONDARY, align=PP_ALIGN.LEFT
        )
        self._add_textbox(
            slide, right_text,
            6.8, y, 5.5, 0.26,
            Theme.FONT_BODY, 11, Theme.TEXT_SECONDARY, align=PP_ALIGN.RIGHT
        )
        self._add_accent_line(slide, Theme.MARGIN_LEFT, 6.92, 11.33, h=0.01, color=Theme.DIVIDER)

    def _add_section_title(self, slide, title, subtitle=None):
        self._add_accent_line(slide, Theme.MARGIN_LEFT, 0.82, 1.05, h=0.05)
        self._add_textbox(
            slide, title,
            Theme.MARGIN_LEFT, 1.0, 9.8, 0.85,
            Theme.FONT_TITLE, 34, Theme.TEXT_PRIMARY, bold=True
        )
        if subtitle:
            self._add_textbox(
                slide, subtitle,
                Theme.MARGIN_LEFT, 1.82, 9.8, 0.45,
                Theme.FONT_BODY, 17, Theme.TEXT_SECONDARY
            )

    def _add_card(self, slide, x, y, w, h, title, body, title_size=18, body_size=15):
        card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = Theme.ACCENT_SOFT
        card.line.color.rgb = RGBColor(203, 218, 242)
        card.line.width = Pt(1)

        self._add_textbox(
            slide, title,
            x + 0.25, y + 0.2, w - 0.5, 0.5,
            Theme.FONT_TITLE, title_size, Theme.TEXT_PRIMARY, bold=True
        )
        self._add_textbox(
            slide, body,
            x + 0.25, y + 0.8, w - 0.5, h - 1.0,
            Theme.FONT_BODY, body_size, Theme.TEXT_SECONDARY, line_spacing=1.3
        )
        return card

    # ---------- Slide builders ----------

    def add_title_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)

        # Vertical accent column
        self._add_accent_line(slide, 0.95, 0.9, 0.06, h=5.7)

        self._add_textbox(
            slide,
            "An Overview of\nFrequency-Modulated\nContinuous-Wave (FMCW) Radar",
            1.3, 1.15, 8.7, 2.7,
            Theme.FONT_TITLE, 45, Theme.TEXT_PRIMARY, bold=True, line_spacing=1.08
        )

        self._add_textbox(
            slide,
            "Course Name / Code",
            1.33, 4.05, 6.2, 0.45,
            Theme.FONT_BODY, 20, Theme.TEXT_SECONDARY
        )

        # Bottom-right metadata block for visual balance
        self._add_card(
            slide,
            8.7, 4.35, 3.65, 1.85,
            "Presenter",
            "Student Name\nUniversity Name",
            title_size=16,
            body_size=14,
        )

        self._add_footer(slide, "Department of Electrical & Computer Engineering", "1")

    def add_outline_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        self._add_section_title(slide, "Outline", "Session roadmap")

        # Two-column structured agenda
        left_x, right_x = Theme.MARGIN_LEFT, 6.95
        top_y = 2.45
        card_w, card_h = 5.35, 1.45

        items = [
            ("01  FMCW Radar Principle", "Linear chirp transmission, beat frequency, and ranging concept."),
            ("02  Signal Chain", "Transmitter, mixer, ADC path, and digital processing overview."),
            ("03  Applications", "Automotive sensing, industrial monitoring, and robotics use cases."),
            ("04  Challenges & Outlook", "Interference, resolution limits, and next-generation trends."),
        ]

        self._add_card(slide, left_x, top_y, card_w, card_h, items[0][0], items[0][1], 16, 14)
        self._add_card(slide, right_x, top_y, card_w, card_h, items[1][0], items[1][1], 16, 14)
        self._add_card(slide, left_x, top_y + card_h + 0.35, card_w, card_h, items[2][0], items[2][1], 16, 14)
        self._add_card(slide, right_x, top_y + card_h + 0.35, card_w, card_h, items[3][0], items[3][1], 16, 14)

        self._add_footer(slide, "An Overview of FMCW Radar", "2")

    def add_conclusion_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        self._add_section_title(slide, "Conclusion", "Key takeaways")

        # Main synthesis block
        self._add_card(
            slide,
            Theme.MARGIN_LEFT,
            2.45,
            11.33,
            2.35,
            "Summary",
            (
                "FMCW radar enables simultaneous range and velocity estimation through chirp-based processing, "
                "offering compact hardware integration and strong potential for robust sensing in dynamic environments."
            ),
            title_size=19,
            body_size=16,
        )

        # Three concise takeaway columns
        box_y = 5.08
        w = 3.55
        gap = 0.34
        starts = [Theme.MARGIN_LEFT, Theme.MARGIN_LEFT + w + gap, Theme.MARGIN_LEFT + 2 * (w + gap)]
        points = [
            ("Core Strength", "High range precision with efficient continuous transmission."),
            ("Practical Value", "Scalable across automotive, biomedical, and smart infrastructure domains."),
            ("Future Direction", "Improved interference mitigation and AI-enhanced signal interpretation."),
        ]

        for x, (t, b) in zip(starts, points):
            self._add_card(slide, x, box_y, w, 1.65, t, b, title_size=15, body_size=13)

        self._add_footer(slide, "An Overview of FMCW Radar", "3")

    def add_thank_you_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)

        # Centered minimal closing composition
        self._add_accent_line(slide, 5.95, 1.72, 1.42, h=0.05)
        self._add_textbox(
            slide,
            "Thank You",
            3.1, 2.05, 7.2, 1.0,
            Theme.FONT_TITLE, 52, Theme.TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER
        )
        self._add_textbox(
            slide,
            "Questions and Discussion",
            3.35, 3.05, 6.7, 0.55,
            Theme.FONT_BODY, 22, Theme.TEXT_SECONDARY, align=PP_ALIGN.CENTER
        )

        # Elegant contact strip
        strip = slide.shapes.add_shape(1, Inches(2.15), Inches(4.15), Inches(9.05), Inches(1.5))
        strip.fill.solid()
        strip.fill.fore_color.rgb = Theme.ACCENT_SOFT
        strip.line.color.rgb = RGBColor(203, 218, 242)
        strip.line.width = Pt(1)

        self._add_textbox(
            slide,
            "Student Name  •  Course Name / Code  •  University Name",
            2.35, 4.7, 8.65, 0.45,
            Theme.FONT_BODY, 16, Theme.TEXT_PRIMARY, align=PP_ALIGN.CENTER
        )

        self._add_footer(slide, "An Overview of FMCW Radar", "4")

    def build(self, output_path="FMCW_Radar_Academic_Template.pptx"):
        self.add_title_slide()
        self.add_outline_slide()
        self.add_conclusion_slide()
        self.add_thank_you_slide()
        self.prs.save(output_path)


if __name__ == "__main__":
    template = AcademicPPTTemplate()
    template.build()
    print("Created FMCW_Radar_Academic_Template.pptx")
